from io import BytesIO
import json

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PRIMARY_DARK = RGBColor(0x0A, 0x16, 0x28)
PRIMARY_BLUE = RGBColor(0x12, 0x1E, 0x36)
ACCENT_RED = RGBColor(0xC6, 0x28, 0x28)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
TABLE_HEADER_BG = RGBColor(0x0A, 0x16, 0x28)
TABLE_ALT_BG = RGBColor(0xF5, 0xF7, 0xFA)


def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size=14,
                 color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_paragraph(text_frame, text, font_size=14, color=DARK_GRAY,
                   bold=False, alignment=PP_ALIGN.LEFT, space_before=0,
                   space_after=0, font_name="Microsoft YaHei"):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p


def _add_table(slide, left, top, width, height, data, col_widths=None):
    rows = len(data)
    cols = len(data[0])
    ts = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for para in cell.text_frame.paragraphs:
                para.font.name = "Microsoft YaHei"
                para.font.size = Pt(11)
                para.alignment = PP_ALIGN.CENTER
                if ri == 0:
                    para.font.color.rgb = WHITE
                    para.font.bold = True
                    para.font.size = Pt(12)
                else:
                    para.font.color.rgb = DARK_GRAY
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return ts


def generate_report_ppt(result):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    emp = result.assignment.employee
    dims = result.dimension_scores
    if isinstance(dims, str):
        dims = json.loads(dims)

    # Slide 1: Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, PRIMARY_DARK)
    _add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_RED)
    _add_textbox(slide, Inches(1), Inches(1.5), Inches(5), Inches(0.6), "\u661f\u6cb3\u667a\u5584\u4eba\u624d\u6d4b\u8bc4", font_size=16, color=MID_GRAY)
    _add_textbox(slide, Inches(1), Inches(2.3), Inches(11), Inches(1.0), "\u5165\u804c\u4eba\u624d\u6d4b\u8bc4\u62a5\u544a", font_size=40, color=WHITE, bold=True)
    _add_shape(slide, Inches(1), Inches(3.5), Inches(3), Inches(0.04), ACCENT_RED)
    _add_textbox(slide, Inches(1), Inches(4.0), Inches(8), Inches(0.5),
                 "\u59d3\u540d\uff1a" + emp.name + "      \u5c97\u4f4d\uff1a" + (emp.position or "\u672a\u586b\u5199"),
                 font_size=18, color=LIGHT_GRAY)
    _add_textbox(slide, Inches(1), Inches(4.6), Inches(8), Inches(0.5),
                 "\u6d4b\u8bc4\u65e5\u671f\uff1a" + result.generated_at.strftime("%Y\u5e74%m\u6708%d\u65e5"),
                 font_size=16, color=MID_GRAY)
    _add_shape(slide, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5), PRIMARY_BLUE)
    _add_textbox(slide, Inches(1), Inches(7.05), Inches(6), Inches(0.4),
                 "Confidential \u00b7 \u661f\u6cb3\u667a\u5584 TIC Talent Assessment System",
                 font_size=9, color=MID_GRAY)

    # Slide 2: Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), PRIMARY_DARK)
    _add_textbox(slide, Inches(0.8), Inches(0.15), Inches(8), Inches(0.6),
                 "\u6d4b\u8bc4\u7efc\u5408\u603b\u89c8", font_size=26, color=WHITE, bold=True)

    risk_labels = {"low": "\u4f4e\u98ce\u9669", "medium": "\u4e2d\u98ce\u9669", "high": "\u9ad8\u98ce\u9669"}
    risk_label = risk_labels.get(result.risk_level, "\u672a\u77e5")
    risk_cm = {"low": RGBColor(0x27, 0xAE, 0x60), "medium": RGBColor(0xF3, 0x9C, 0x12), "high": ACCENT_RED}
    risk_color = risk_cm.get(result.risk_level, MID_GRAY)

    cl, ct, cw, ch = Inches(0.8), Inches(1.4), Inches(3.8), Inches(2.2)
    _add_shape(slide, cl, ct, cw, ch, LIGHT_GRAY)
    _add_textbox(slide, cl+Inches(0.3), ct+Inches(0.3), cw-Inches(0.6), Inches(0.4),
                 "\u98ce\u9669\u8bc4\u4f30", font_size=13, color=MID_GRAY)
    _add_textbox(slide, cl+Inches(0.3), ct+Inches(0.9), cw-Inches(0.6), Inches(0.8),
                 risk_label, font_size=36, color=risk_color, bold=True)

    cm = Inches(4.9)
    _add_shape(slide, cm, ct, cw, ch, LIGHT_GRAY)
    _add_textbox(slide, cm+Inches(0.3), ct+Inches(0.3), cw-Inches(0.6), Inches(0.4),
                 "\u7efc\u5408\u5f97\u5206\u7387", font_size=13, color=MID_GRAY)
    _add_textbox(slide, cm+Inches(0.3), ct+Inches(0.9), cw-Inches(0.6), Inches(0.8),
                 "%.1f%%" % result.score_percent, font_size=36, color=PRIMARY_DARK, bold=True)

    cr = Inches(9.0)
    _add_shape(slide, cr, ct, cw, ch, LIGHT_GRAY)
    _add_textbox(slide, cr+Inches(0.3), ct+Inches(0.3), cw-Inches(0.6), Inches(0.4),
                 "\u5c97\u4f4d\u9002\u914d\u5ea6", font_size=13, color=MID_GRAY)
    _add_textbox(slide, cr+Inches(0.3), ct+Inches(0.9), cw-Inches(0.6), Inches(0.8),
                 "%.1f" % result.fit_score, font_size=36, color=ACCENT_RED, bold=True)

    ty = Inches(4.0)
    if result.risk_tags:
        _add_textbox(slide, Inches(0.8), ty, Inches(3), Inches(0.4),
                     "\u91cd\u70b9\u5173\u6ce8\u9879", font_size=15, color=DARK_GRAY, bold=True)
        _add_textbox(slide, Inches(0.8), ty+Inches(0.5), Inches(11), Inches(0.4),
                     "\u3001".join(result.risk_tags), font_size=13, color=ACCENT_RED)
    if result.is_abnormal:
        ay = ty + Inches(1.2) if result.risk_tags else ty
        _add_textbox(slide, Inches(0.8), ay, Inches(11), Inches(0.4),
                     "\u26a0 \u6ce8\u610f\uff1a\u8be5\u6d4b\u8bc4\u53ef\u80fd\u5b58\u5728\u5f02\u5e38\uff08" + result.abnormal_reason + "\uff09\uff0c\u5efa\u8baeHR\u590d\u6838",
                     font_size=12, color=ACCENT_RED)

    iy = Inches(5.6)
    _add_shape(slide, Inches(0.8), iy, Inches(11.5), Inches(1.2), LIGHT_GRAY)
    gs = "\u7537" if emp.gender == "male" else "\u5973" if emp.gender == "female" else "\u672a\u586b"
    ln1 = "\u59d3\u540d\uff1a%s    \u6027\u522b\uff1a%s    \u5e74\u9f84\uff1a%s" % (emp.name, gs, emp.age or "-")
    ln2 = "\u90e8\u95e8\uff1a%s    \u5c97\u4f4d\uff1a%s    \u6240\u5c5e\u5206\u516c\u53f8\uff1a%s" % (emp.department or "-", emp.position or "-", emp.branch.name if emp.branch else "-")
    tb = _add_textbox(slide, Inches(1.0), iy+Inches(0.1), Inches(11), Inches(1.0), ln1, font_size=11, color=DARK_GRAY)
    _add_paragraph(tb.text_frame, ln2, font_size=11, color=DARK_GRAY)

    # Slide 3: Dimension Scores
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), PRIMARY_DARK)
    _add_textbox(slide, Inches(0.8), Inches(0.15), Inches(8), Inches(0.6),
                 "\u5404\u7ef4\u5ea6\u5f97\u5206\u8be6\u60c5", font_size=26, color=WHITE, bold=True)

    if dims:
        td = [["\u7ef4\u5ea6", "\u5f97\u5206", "\u6ee1\u5206", "\u5f97\u5206\u7387", "\u8bc4\u4f30\u7b49\u7ea7"]]
        for d in dims.values():
            td.append([d.get("name", ""), "%.1f" % d.get("score", 0),
                       "%.1f" % d.get("max", 0), "%.1f%%" % d.get("percent", 0),
                       d.get("level", "-")])
        _add_table(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.4 * len(td)),
                   td, col_widths=[Inches(3.5), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)])

        yp = Inches(1.3) + Inches(0.4 * len(td)) + Inches(0.4)
        _add_textbox(slide, Inches(0.8), yp, Inches(11), Inches(0.4),
                     "\u8bc4\u4f30\u7b49\u7ea7\u8bf4\u660e", font_size=14, color=DARK_GRAY, bold=True)
        lg = ("\u25cf \u826f\u597d\uff08\u5f97\u5206\u7387 \u2265 60%\uff09\uff1a\u8be5\u7ef4\u5ea6\u8868\u73b0\u6b63\u5e38\uff0c\u5177\u5907\u76f8\u5e94\u80fd\u529b\u7d20\u8d28\u3002\n"
              "\u25cf \u4e2d\u98ce\u9669\uff0840% \u2264 \u5f97\u5206\u7387 < 60%\uff09\uff1a\u8be5\u7ef4\u5ea6\u9700\u5173\u6ce8\uff0c\u5efa\u8bae\u4e86\u89e3\u5177\u4f53\u60c5\u51b5\u3002\n"
              "\u25cf \u9ad8\u98ce\u9669\uff08\u5f97\u5206\u7387 < 40%\uff09\uff1a\u8be5\u7ef4\u5ea6\u9700\u91cd\u70b9\u5173\u6ce8\uff0c\u5efa\u8bae\u8fdb\u4e00\u6b65\u8bc4\u4f30\u3002")
        _add_textbox(slide, Inches(0.8), yp+Inches(0.4), Inches(11), Inches(1.2), lg, font_size=11, color=MID_GRAY)

    # Slide 4: Suggestions
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), PRIMARY_DARK)
    _add_textbox(slide, Inches(0.8), Inches(0.15), Inches(8), Inches(0.6),
                 "\u7efc\u5408\u5efa\u8bae\u4e0eHR\u5907\u6ce8", font_size=26, color=WHITE, bold=True)

    if result.suggestion:
        ss = result.suggestion.replace("\r\n", "\n").split("\n")
        _add_textbox(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.4),
                     "\u7cfb\u7edf\u7efc\u5408\u5efa\u8bae", font_size=15, color=DARK_GRAY, bold=True)
        st = "\n".join("  \u2022 " + s.strip() for s in ss if s.strip())
        _add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(2.5), st, font_size=12, color=DARK_GRAY)
    else:
        _add_textbox(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.4),
                     "\u7cfb\u7edf\u7efc\u5408\u5efa\u8bae", font_size=15, color=DARK_GRAY, bold=True)
        _add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.4),
                     "\u6682\u65e0\u7cfb\u7edf\u7efc\u5408\u5efa\u8bae\u3002", font_size=12, color=MID_GRAY)

    if result.hr_comment:
        _add_textbox(slide, Inches(0.8), Inches(4.0), Inches(3), Inches(0.4),
                     "HR \u5907\u6ce8", font_size=15, color=DARK_GRAY, bold=True)
        _add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.5),
                     result.hr_comment, font_size=12, color=DARK_GRAY)

    # Slide 5: End
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, PRIMARY_DARK)
    _add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_RED)
    _add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.0),
                 "\u62a5\u544a\u5b8c\u6bd5", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    _add_shape(slide, Inches(5.5), Inches(3.7), Inches(2.3), Inches(0.04), ACCENT_RED)
    _add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
                 "\u611f\u8c22\u53c2\u4e0e\u672c\u6b21\u4eba\u624d\u8bc4\u4f30", font_size=18, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    _add_shape(slide, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5), PRIMARY_BLUE)
    _add_textbox(slide, Inches(1), Inches(7.05), Inches(6), Inches(0.4),
                 "Confidential \u00b7 \u661f\u6cb3\u667a\u5584 TIC Talent Assessment System",
                 font_size=9, color=MID_GRAY)

    buf = BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()
    buf.close()
    return pptx_bytes
